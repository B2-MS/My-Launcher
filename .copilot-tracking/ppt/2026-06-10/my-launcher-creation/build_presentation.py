from pathlib import Path
import struct
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path('/Users/barticus/VS Code/My Launcher')
WORKDIR = ROOT / '.copilot-tracking' / 'ppt' / '2026-06-10' / 'my-launcher-creation'
OUTPUT = WORKDIR / 'slide-deck' / 'My Launcher - Creation Story.pptx'
IMAGES = ROOT / 'docs' / 'images'

BG = '0B0D10'
PANEL = '14181F'
PANEL_2 = '1B212A'
PANEL_EDGE = '2D3642'
ACCENT = '0050EF'
ACCENT_2 = '4B9BFF'
TEXT = 'F5F7FA'
MUTED = 'A6AFBD'
SUBTLE = '6F7A8A'
WHITE = 'FFFFFF'


def hex_to_rgb(value):
    value = value.lstrip('#')
    return tuple(int(value[index:index + 2], 16) for index in (0, 2, 4))


def rgb(value):
    return RGBColor(*hex_to_rgb(value))


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_shape(slide, shape_type, left, top, width, height, fill=None, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=20,
    font_name='Segoe UI',
    color=TEXT,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    fill=None,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(fill)
    else:
        box.fill.background()
    box.line.fill.background()
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.clear()

    lines = text.split('\n')
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color)
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
    return box


def add_pill(slide, left, top, width, height, text, fill=ACCENT, font_size=12, color=WHITE):
    pill = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill, line=fill, radius=True)
    add_textbox(slide, left, top + 0.02, width, height - 0.04, text, font_size=font_size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return pill


def add_header(slide, title, subtitle=None, section=None, slide_num=None):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 0.12, fill=ACCENT, line=ACCENT)
    if section:
        add_textbox(slide, 0.75, 0.38, 5.6, 0.25, section.upper(), font_size=11, color=ACCENT_2, bold=True)
    add_textbox(slide, 0.75, 0.58, 10.8, 0.55, title, font_size=26, color=TEXT, bold=True)
    if subtitle:
        add_textbox(slide, 0.75, 1.05, 11.0, 0.45, subtitle, font_size=12, color=MUTED)
    if slide_num is not None:
        add_textbox(slide, 12.3, 7.0, 0.6, 0.25, str(slide_num), font_size=11, color=SUBTLE, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, bullets, fill=PANEL, line=PANEL_EDGE, accent=ACCENT):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill, line=line, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, 0.08, height, fill=accent, line=accent)
    add_textbox(slide, left + 0.22, top + 0.16, width - 0.36, 0.32, title, font_size=14, color=WHITE, bold=True)
    body = '\n'.join(f'- {bullet}' for bullet in bullets)
    add_textbox(slide, left + 0.22, top + 0.56, width - 0.38, height - 0.72, body, font_size=12, color=MUTED)


def image_size(path):
    with open(path, 'rb') as handle:
        header = handle.read(24)
        if header.startswith(b'\x89PNG\r\n\x1a\n'):
            width, height = struct.unpack('>II', header[16:24])
            return width, height
        raise ValueError(f'Unsupported image format for {path}')


def add_image_fit(slide, path, left, top, width, height, border=PANEL_EDGE, fill=PANEL):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill, line=border, radius=True)
    img_w, img_h = image_size(path)
    aspect = img_w / img_h
    padding = 0.14
    inner_left = left + padding
    inner_top = top + padding
    inner_width = width - padding * 2
    inner_height = height - padding * 2
    inner_aspect = inner_width / inner_height
    if aspect > inner_aspect:
        pic_width = inner_width
        pic_height = inner_width / aspect
        pic_left = inner_left
        pic_top = inner_top + (inner_height - pic_height) / 2
    else:
        pic_height = inner_height
        pic_width = inner_height * aspect
        pic_left = inner_left + (inner_width - pic_width) / 2
        pic_top = inner_top
    slide.shapes.add_picture(str(path), Inches(pic_left), Inches(pic_top), Inches(pic_width), Inches(pic_height))


def add_callout(slide, left, top, width, height, title, body, accent=ACCENT):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=PANEL_2, line=PANEL_EDGE, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, 0.08, height, fill=accent, line=accent)
    add_textbox(slide, left + 0.2, top + 0.14, width - 0.3, 0.25, title, font_size=12, color=WHITE, bold=True)
    add_textbox(slide, left + 0.2, top + 0.43, width - 0.32, height - 0.52, body, font_size=11, color=MUTED)


def add_connector(slide, x1, y1, x2, y2, color=ACCENT_2, width=2):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    return connector


def set_notes(slide, text):
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.7, 1.0, 6.1, 5.3, fill=PANEL, line=PANEL_EDGE, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.7, 1.0, 0.1, 5.3, fill=ACCENT, line=ACCENT)
    add_textbox(slide, 1.0, 1.25, 4.9, 0.35, 'CASE STUDY', font_size=12, color=ACCENT_2, bold=True)
    add_textbox(slide, 1.0, 1.62, 5.0, 1.35, 'My Launcher', font_size=30, color=WHITE, bold=True)
    add_textbox(slide, 1.0, 3.02, 5.0, 0.65, 'How the app was created', font_size=19, color=MUTED)
    add_textbox(slide, 1.0, 3.78, 5.2, 0.7, 'Jetpack Compose launcher inspired by Windows Phone Metro UI', font_size=16, color=TEXT)
    add_pill(slide, 1.0, 4.75, 1.55, 0.36, 'Compose')
    add_pill(slide, 2.72, 4.75, 0.95, 0.36, 'Hilt', fill='0F5BB5')
    add_pill(slide, 3.8, 4.75, 1.4, 0.36, 'DataStore', fill='1E6DD1')
    add_pill(slide, 5.35, 4.75, 1.2, 0.36, 'Launcher', fill='2D3642')
    add_textbox(slide, 1.0, 5.42, 4.8, 0.45, 'Alpha build for Android 10+ with a 6-column tile grid, edit mode, app list, and customization screens.', font_size=12, color=MUTED)
    add_image_fit(slide, IMAGES / 'start-screen.png', 7.2, 0.75, 5.3, 6.1)
    add_textbox(slide, 7.45, 6.95, 4.8, 0.22, 'Start screen preview from the running app', font_size=10, color=SUBTLE)
    set_notes(slide, 'Open with the core story: the deck explains how the launcher was planned from the PRD, built in Compose, and validated on an Android emulator. Emphasize that the app targets a Windows Phone-style experience rather than a generic Android home screen.')
    return slide


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Why this app exists', 'The idea came from a gap in the Android launcher market.', 'Problem and motivation', 2)
    add_card(slide, 0.8, 1.55, 5.75, 2.2, 'The gap', [
        'Windows Phone Metro UI is still distinctive and remembered fondly.',
        'Most Android launchers copy the same icon-grid pattern.',
        'Existing Metro-style launchers are usually abandoned or inaccurate.'
    ])
    add_card(slide, 0.8, 4.0, 5.75, 2.2, 'The goal', [
        'Build a launcher that feels faithful, modern, and maintainable.',
        'Support real customization: tile size, transparency, colors, and layout.',
        'Use modern Android architecture instead of a one-off UI demo.'
    ], fill='111720')
    add_card(slide, 6.85, 1.55, 5.65, 4.7, 'Product intent', [
        'A home-screen replacement registered as HOME + LAUNCHER.',
        'A vertically scrolling 6-column Start screen with resizable tiles.',
        'An alphabetical app list with search and pin-to-start flows.',
        'A settings surface for accent color, glass bevel, and opacity.',
        'A path toward live tiles, persistence, and future widget integration.'
    ], fill='111720')
    set_notes(slide, 'Frame the product as a response to a real UX gap. The PRD described former Windows Phone users, customization enthusiasts, and minimalist users as the target audience. The implementation was designed to satisfy all three with a clean launcher that still behaves like Android.')
    return slide


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'From PRD to alpha', 'The build sequence was intentionally narrow: define the core experience first, then polish it.', 'Implementation path', 3)
    steps = [
        ('1', 'Translate the PRD into a small alpha scope', 'Focus on the Start screen, app list, edit mode, and settings.'),
        ('2', 'Choose a modern stack', 'Kotlin, Jetpack Compose, Hilt, DataStore, Room, and Gradle Kotlin DSL.'),
        ('3', 'Build the launcher shell', 'Register the activity as HOME + LAUNCHER and add edge-to-edge layout.'),
        ('4', 'Implement the tile engine', 'Layout tiles on a 6-column grid with resize and drag behavior.'),
        ('5', 'Validate on-device', 'Install to an emulator, inspect screenshots, and iterate on spacing and behavior.'),
    ]
    y = 1.55
    for number, title, body in steps:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, y, 11.65, 0.82, fill=PANEL, line=PANEL_EDGE, radius=True)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 1.06, y + 0.15, 0.5, 0.5, fill=ACCENT, line=ACCENT)
        add_textbox(slide, 1.06, y + 0.18, 0.5, 0.3, number, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 1.78, y + 0.1, 4.8, 0.2, title, font_size=14, color=WHITE, bold=True)
        add_textbox(slide, 1.78, y + 0.35, 9.9, 0.25, body, font_size=11, color=MUTED)
        y += 0.95
    add_callout(slide, 8.1, 6.0, 4.3, 0.78, 'Alpha scope', 'The first release focused on a believable launcher loop instead of every Metro feature from the PRD.')
    set_notes(slide, 'Explain that the first version was intentionally scoped to the highest-value launcher behaviors. The alpha build proves the shell, tile grid, app discovery, settings, and edit loop before expanding into live content or sync features.')
    return slide


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Architecture', 'The app is organized around a simple flow: discover apps, store tile state, and render it in Compose.', 'Data flow and control surfaces', 4)
    boxes = [
        (0.9, 1.95, 2.2, 0.75, 'Installed apps', 'PackageManager / LauncherApps'),
        (3.45, 1.95, 2.2, 0.75, 'Repositories', 'App discovery and tile data'),
        (5.98, 1.95, 2.2, 0.75, 'ViewModel', 'Single source of state'),
        (8.5, 1.95, 2.2, 0.75, 'Compose UI', 'Start screen, list, settings'),
    ]
    for left, top, width, height, title, body in boxes:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=PANEL, line=PANEL_EDGE, radius=True)
        add_textbox(slide, left + 0.12, top + 0.08, width - 0.24, 0.18, title, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + 0.12, top + 0.34, width - 0.24, 0.16, body, font_size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_connector(slide, 3.15, 2.33, 3.45, 2.33)
    add_connector(slide, 5.68, 2.33, 5.98, 2.33)
    add_connector(slide, 8.2, 2.33, 8.5, 2.33)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.05, 3.35, 4.35, 1.22, fill='10161D', line=PANEL_EDGE, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.65, 3.35, 4.65, 1.22, fill='10161D', line=PANEL_EDGE, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.05, 4.82, 9.25, 1.22, fill='10161D', line=PANEL_EDGE, radius=True)
    add_textbox(slide, 1.28, 3.55, 3.8, 0.2, 'User actions', font_size=13, color=WHITE, bold=True)
    add_textbox(slide, 1.28, 3.86, 3.9, 0.5, 'Tap, long-press, drag, resize, theme changes, and group edits all feed back into the ViewModel.', font_size=11, color=MUTED)
    add_textbox(slide, 5.88, 3.55, 4.0, 0.2, 'Persistence', font_size=13, color=WHITE, bold=True)
    add_textbox(slide, 5.88, 3.86, 4.0, 0.5, 'Room is declared for tile storage while DataStore handles preferences such as theme and transparency.', font_size=11, color=MUTED)
    add_textbox(slide, 1.28, 5.02, 8.75, 0.2, 'Manifest and platform integration', font_size=13, color=WHITE, bold=True)
    add_textbox(slide, 1.28, 5.34, 8.75, 0.45, 'The activity is registered as both MAIN/LAUNCHER and MAIN/HOME, which is what turns the app into a real home-screen replacement.', font_size=11, color=MUTED)
    set_notes(slide, 'This is the technical core slide. The launcher works because app discovery, state management, persistence, and Compose rendering are kept in separate layers. Emphasize that the activity is registered as a home app, not just a normal screen in a normal app.')
    return slide


def build_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Stack and tooling', 'The implementation uses current Android tooling, but stays conservative where stability matters.', 'Build system and dependencies', 5)
    add_card(slide, 0.8, 1.55, 3.95, 2.25, 'UI layer', [
        'Kotlin + Jetpack Compose',
        'Navigation Compose and activity-compose',
        'Compose animations for tile transitions and press tilt',
        'Material 3 for the settings surface'
    ], fill='111720')
    add_card(slide, 4.75, 1.55, 3.95, 2.25, 'State and data', [
        'Hilt for dependency injection',
        'DataStore preferences for user settings',
        'Room for tile layout persistence',
        'Repository pattern for app discovery and launch data'
    ], fill='111720')
    add_card(slide, 8.7, 1.55, 3.85, 2.25, 'Build and release', [
        'Gradle Kotlin DSL with version catalogs',
        'Android Studio JBR for builds and installs',
        'Min SDK 29, target SDK 35',
        'Debug installs via installDebug on emulator'
    ], fill='111720')
    add_card(slide, 0.8, 4.15, 11.75, 1.7, 'Why this stack fits the app', [
        'Compose makes the custom tile UI and drag/edit interactions practical.',
        'Hilt and repositories keep the launcher code testable as it grows.',
        'DataStore and Room separate small preferences from structured layout state.'
    ])
    set_notes(slide, 'Call out that the stack is intentionally modern but not experimental. Compose handles the launcher UI well, while DataStore and Room keep settings and tile state separate from rendering logic. The version catalog keeps dependency management tidy.')
    return slide


def build_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Start screen', 'This is the main home screen: a six-column tile grid with drag-and-drop behavior.', 'Metro-style home shell', 6)
    add_image_fit(slide, IMAGES / 'start-screen.png', 0.85, 1.55, 5.65, 5.65)
    add_callout(slide, 6.75, 1.65, 5.7, 1.08, 'Grid behavior', 'Tiles can span 1 to 6 columns wide and 1 to 4 rows tall, which keeps the layout flexible while preserving the Metro rhythm.')
    add_callout(slide, 6.75, 2.92, 5.7, 1.08, 'Interaction', 'Long-press enters edit mode; dragging automatically displaces overlapping tiles to the nearest free space.')
    add_callout(slide, 6.75, 4.19, 5.7, 1.08, 'Visual treatment', 'Wallpaper shows through transparent tiles, and the UI uses large type, dark surfaces, and subtle glass bevels.')
    add_callout(slide, 6.75, 5.46, 5.7, 1.08, 'Start screen loop', 'This screen is the reason the app feels like a launcher instead of a standard list-based application.')
    set_notes(slide, 'Use the screenshot to show the actual result of the tile grid implementation. Mention that the UI is designed around a six-column grid, with the wallpaper visible through translucent tiles and a start screen layout that mirrors Windows Phone.')
    return slide


def build_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Edit mode and groups', 'The tile system supports resizing, unpinning, and dwell-to-group interactions.', 'Tile management', 7)
    add_image_fit(slide, IMAGES / 'edit-mode.png', 0.75, 1.55, 5.85, 5.55)
    add_image_fit(slide, IMAGES / 'tile-group-expanded.png', 6.75, 1.55, 5.85, 2.55)
    add_callout(slide, 6.75, 4.3, 5.85, 0.86, 'Edit mode', 'Long-press a tile to shrink it, show its size label, and expose the unpin button.')
    add_callout(slide, 6.75, 5.28, 5.85, 0.86, 'Grouping', 'Dragging onto another tile and holding still creates a group; expanded groups can be rearranged and ungrouped.')
    set_notes(slide, 'This slide explains the more distinctive launcher behavior. Edit mode exposes the controls that make the layout personal, while dwell-to-group is the main interaction that differentiates the app from a standard icon grid.')
    return slide


def build_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'App list and search', 'The second screen gives fast discovery and a clean pin-to-start workflow.', 'Navigation and discovery', 8)
    add_image_fit(slide, IMAGES / 'app-list.png', 0.85, 1.55, 5.15, 5.7)
    add_card(slide, 6.35, 1.65, 6.05, 1.25, 'Navigation model', [
        'Swipe right from Start, or tap All Apps at the bottom-right.',
        'The app list uses alphabetical headers and a search bar at the top.'
    ], fill='111720')
    add_card(slide, 6.35, 3.05, 6.05, 1.25, 'Pinning flow', [
        'Long-press any app to pin it back to the Start screen as a tile.',
        'This keeps the launcher tied to the installed app set on the device.'
    ], fill='111720')
    add_card(slide, 6.35, 4.45, 6.05, 1.25, 'Visual pattern', [
        'Accent-colored square icon blocks echo the Windows Phone app list.',
        'The layout is intentionally dense and text-forward.'
    ], fill='111720')
    set_notes(slide, 'Explain that the app list is not a separate launcher, but the second half of the same navigation model. It supports quick search, section headers, and a long-press pin action that feeds the Start screen.')
    return slide


def build_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Settings and theming', 'The customization surface makes the launcher feel personal instead of fixed.', 'Visual customization', 9)
    add_image_fit(slide, IMAGES / 'settings-screen.png', 0.85, 1.55, 4.85, 5.7)
    add_card(slide, 6.1, 1.65, 6.3, 1.0, 'Accent color', [
        'HSV picker with live hex preview and a Metro-style blue default.'
    ], fill='111720')
    add_card(slide, 6.1, 2.82, 6.3, 1.0, 'Transparency and bevel', [
        'Global tile opacity slider plus a glass-like bevel depth control.'
    ], fill='111720')
    add_card(slide, 6.1, 3.99, 6.3, 1.0, 'Theme and interval', [
        'Dark/light mode, live-tile animation interval, and theme saving.'
    ], fill='111720')
    add_card(slide, 6.1, 5.16, 6.3, 1.0, 'Why settings matter', [
        'The UI design only works if users can shape the look and density of the screen.'
    ], fill='111720')
    set_notes(slide, 'Show the settings screen as the control center for the app. The important detail is that the launcher is highly themeable, so the same core UI can support different tastes and levels of transparency.')
    return slide


def build_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Build and validation loop', 'The app was created by building, installing, and checking the result on a live emulator.', 'Development workflow', 10)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.95, 2.0, 2.0, 1.0, fill=PANEL, line=PANEL_EDGE, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 3.35, 2.0, 2.0, 1.0, fill=PANEL, line=PANEL_EDGE, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.75, 2.0, 2.0, 1.0, fill=PANEL, line=PANEL_EDGE, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 8.15, 2.0, 2.0, 1.0, fill=PANEL, line=PANEL_EDGE, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.55, 2.0, 2.0, 1.0, fill=PANEL, line=PANEL_EDGE, radius=True)
    add_textbox(slide, 1.12, 2.22, 1.65, 0.3, 'Gradle', font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 3.52, 2.22, 1.65, 0.3, 'Emulator', font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 5.92, 2.22, 1.65, 0.3, 'adb install', font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 8.32, 2.22, 1.65, 0.3, 'Launch', font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 10.72, 2.22, 1.65, 0.3, 'Inspect', font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_connector(slide, 2.95, 2.5, 3.35, 2.5)
    add_connector(slide, 5.35, 2.5, 5.75, 2.5)
    add_connector(slide, 7.75, 2.5, 8.15, 2.5)
    add_connector(slide, 10.15, 2.5, 10.55, 2.5)
    add_callout(slide, 0.95, 4.1, 3.1, 1.3, 'Environment', 'Android Studio JBR supplies the JDK, while the local Android SDK provides emulator and platform-tools.')
    add_callout(slide, 4.2, 4.1, 3.1, 1.3, 'Validation', 'The debug build is installed with installDebug, then the launcher is opened directly on the device.')
    add_callout(slide, 7.45, 4.1, 5.1, 1.3, 'Documentation', 'Screenshots and release notes in docs/ capture the state of the app and the roadmap after each milestone.')
    set_notes(slide, 'This is the engineering loop slide. The project was validated by building to a running emulator, launching the home activity, and comparing the result to the documented UI. That keeps the deck grounded in an actual implementation path rather than a slide-only concept.')
    return slide


def build_slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Where it stands now', 'The alpha is a working launcher shell with a clear path for the next features.', 'Status and next steps', 11)
    add_card(slide, 0.8, 1.6, 3.8, 4.8, 'What is in place', [
        'Working home-screen replacement',
        'Start screen tile grid and app list',
        'Edit mode, drag and drop, and groups',
        'Accent color and transparency settings',
        'Metro-inspired visual language'
    ], fill='111720')
    add_card(slide, 4.8, 1.6, 3.8, 4.8, 'What is still evolving', [
        'Persistent tile storage refinement',
        'True live tile data sources',
        'Notification badges',
        'Backup and restore polish',
        'Widget and richer integration work'
    ], fill='111720')
    add_card(slide, 8.8, 1.6, 3.7, 4.8, 'Takeaway', [
        'The project is already a usable launcher.',
        'The remaining work is about fidelity, data, and scale.',
        'The architecture leaves room for those additions without a rewrite.'
    ], fill='111720')
    add_textbox(slide, 0.8, 6.65, 11.7, 0.35, 'Result: a modern Android launcher that recreates the Windows Phone feel while staying maintainable in Kotlin and Compose.', font_size=13, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    set_notes(slide, 'Close with the current state of the app. The alpha already demonstrates the key user experience, and the remaining roadmap is about making the launcher richer and more durable over time.')
    return slide


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        build_slide_1,
        build_slide_2,
        build_slide_3,
        build_slide_4,
        build_slide_5,
        build_slide_6,
        build_slide_7,
        build_slide_8,
        build_slide_9,
        build_slide_10,
        build_slide_11,
    ]
    for builder in builders:
        builder(prs)
    prs.save(str(OUTPUT))
    print(f'Wrote {OUTPUT}')


if __name__ == '__main__':
    main()
